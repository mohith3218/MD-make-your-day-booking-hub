import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Deep Modern Tech
    BG_DARK = RGBColor(15, 23, 42)       # Slate 900
    CARD_BG = RGBColor(30, 41, 59)       # Slate 800
    CARD_BORDER = RGBColor(51, 65, 85)   # Slate 700
    PRIMARY = RGBColor(99, 102, 241)     # Indigo 500
    ACCENT = RGBColor(236, 72, 153)      # Pink/Rose 500
    CYAN = RGBColor(6, 182, 212)         # Cyan 500
    TEXT_LIGHT = RGBColor(248, 250, 252) # Slate 50
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
    SUCCESS = RGBColor(16, 185, 129)     # Emerald 500

    def add_blank_slide():
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        # Background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return slide

    def add_header(slide, title_text, subtitle_text=None, category="CINEWAVE • 2ND REVIEW"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY

        # Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.6))
        tf = t_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT

        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.4))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide1 = add_blank_slide()
    # Glow / Accent Card in center
    add_card(slide1, Inches(1.2), Inches(1.2), Inches(10.9), Inches(5.1), CARD_BG, PRIMARY)

    # Title text
    tb = slide1.shapes.add_textbox(Inches(1.8), Inches(1.6), Inches(9.7), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "PYTHON FULL STACK PROJECT • 2ND MILESTONE REVIEW"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = CYAN
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "CineWave: Movie Booking & Cinema Management Hub"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "A full-stack, concurrency-safe cinema platform featuring custom auditorium seat layouts, dual payment workflows, dynamic pricing, and automated management."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "Presented by: Python Full Stack Student  |  Tech Stack: Python (Flask) • SQLite • REST APIs • Razorpay"
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = PRIMARY

    # ==================== SLIDE 2: PROJECT OBJECTIVES & SCOPE ====================
    slide2 = add_blank_slide()
    add_header(slide2, "Project Objectives & 2nd Review Scope", "Bridging the gap between static booking sites and enterprise-grade real-time systems")

    cards_data = [
        ("Problem Statement", "Existing academic booking apps use static seat numbers without real hall geometry, lack concurrency-safe seat locking, and use dummy unverified checkout flows.", ACCENT),
        ("2nd Review Goals", "Deliver an end-to-end operational prototype: dynamic JSON seating maps, 8-minute distributed seat locks, dual payment gateway (Razorpay + Sandbox), and admin control.", CYAN),
        ("Core Deliverables", "• Realistic multi-tier theater seat layouts\n• Double-booking prevention algorithm\n• Full payment lifecycle with signature checks\n• In-app wallet ledger & admin recharges\n• Dynamic demand-based surge pricing", SUCCESS),
    ]

    for i, (title, desc, color) in enumerate(cards_data):
        c_left = Inches(0.8 + i * 3.95)
        add_card(slide2, c_left, Inches(2.0), Inches(3.75), Inches(4.7))
        tb = slide2.shapes.add_textbox(c_left + Inches(0.25), Inches(2.3), Inches(3.25), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.space_after = Pt(12)

        for line in desc.split("\n"):
            p_b = tf.add_paragraph()
            p_b.text = line
            p_b.font.size = Pt(13)
            p_b.font.color.rgb = TEXT_LIGHT
            p_b.space_after = Pt(8)

    # ==================== SLIDE 3: SYSTEM ARCHITECTURE & TECH STACK ====================
    slide3 = add_blank_slide()
    add_header(slide3, "Full-Stack System Architecture & Tech Stack", "Modular MVC design built on Python Flask and modern web standards")

    layers = [
        ("Frontend Layer", "Jinja2 Templates • HTML5 • Modern CSS3 • JavaScript (Fetch API polling for live seat updates & dynamic totals) • QR Code generation"),
        ("Application Layer", "Flask Framework (Python 3.x) • RESTful API Endpoints • Session Management • Werkzeug Security • Seating Geometry Engine (seat_layout.py)"),
        ("Payment & Services", "Razorpay Checkout API • Mock Gateway Sandbox • HMAC-SHA256 Signature Verification • TMDb Movie API Integration"),
        ("Database & ORM Layer", "SQLite with SQLAlchemy ORM • Relational Integrity with Foreign Keys & Unique Constraints • Auto-migration engine")
    ]

    for i, (title, desc) in enumerate(layers):
        c_top = Inches(1.9 + i * 1.25)
        add_card(slide3, Inches(0.8), c_top, Inches(11.7), Inches(1.1))
        tb = slide3.shapes.add_textbox(Inches(1.1), c_top + Inches(0.15), Inches(11.1), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title + "  ➔  "
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_LIGHT

    # ==================== SLIDE 4: DATABASE SCHEMA & ENTITY MODEL ====================
    slide4 = add_blank_slide()
    add_header(slide4, "Database Design & Entity Relational Schema", "8 Normalized relational models ensuring complete data integrity")

    # Table Left: Core Entities
    add_card(slide4, Inches(0.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb1 = slide4.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "Cinema & Catalog Entities"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(10)

    entities_left = [
        ("User", "id, username, email, password_hash, points, is_admin"),
        ("Movie", "id, title, description, duration, poster, language, certificate"),
        ("Theater", "id, name, location, city, address (1-to-many with Screens)"),
        ("Screen", "id, theater_id, name, screen_type, layout_json"),
        ("Showtime", "id, movie_id, screen_id, show_date, show_time, price_overrides_json")
    ]
    for name, fields in entities_left:
        p = tf1.add_paragraph()
        p.text = f"• {name}: {fields}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(6)

    # Table Right: Booking & Financial Entities
    add_card(slide4, Inches(6.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb2 = slide4.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.3), Inches(4.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Transactional & Locking Entities"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.space_after = Pt(10)

    entities_right = [
        ("SeatLock", "id, showtime_id, user_id, token, seats, expires_at (8-min TTL)"),
        ("Booking", "id, reference, user_id, showtime_id, seats, total_amount, status"),
        ("BookedSeat", "id, showtime_id, seat_id, booking_id [UNIQUE(showtime_id, seat_id)]"),
        ("Payment", "id, booking_id, gateway, order_id, payment_id, signature, status"),
        ("Wallet & Ledger", "Wallet, WalletTransaction, RechargeRequest (Audit history)")
    ]
    for name, fields in entities_right:
        p = tf2.add_paragraph()
        p.text = f"• {name}: {fields}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(6)

    # ==================== SLIDE 5: DYNAMIC SEATING ENGINE ====================
    slide5 = add_blank_slide()
    add_header(slide5, "Core Innovation: Dynamic Theater Seating Engine", "Theater ➔ Screen ➔ JSON Layout architecture replicating authentic cinema auditoriums")

    # Left: Concept
    add_card(slide5, Inches(0.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb_s5_1 = slide5.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf_s5_1 = tb_s5_1.text_frame
    tf_s5_1.word_wrap = True

    p = tf_s5_1.paragraphs[0]
    p.text = "Geometry & Layout Capabilities"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)

    feats = [
        "Multi-Class Tiering: Supports Recliner, Prime, Classic, Executive tiers per screen with distinct pricing.",
        "Aisles & Walkways: Configurable aisles_after positions to create realistic cinema gangways.",
        "Physical Anomalies: Supports blocked non-existent architectural seats and offset for centered curved rows.",
        "Seat Spans: Double-width couple/sofa seats (seat_span: 2).",
        "Visual Admin Editor: Real-time interactive layout preview with JSON validator."
    ]
    for f in feats:
        p = tf_s5_1.add_paragraph()
        p.text = "• " + f
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

    # Right: Seeded Theatres Data
    add_card(slide5, Inches(6.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb_s5_2 = slide5.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf_s5_2 = tb_s5_2.text_frame
    tf_s5_2.word_wrap = True

    p = tf_s5_2.paragraphs[0]
    p.text = "Real-World Benchmark (Visakhapatnam)"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.space_after = Pt(10)

    theaters_list = [
        "23 Theatres, 53 Screens, 15,149 Total Seats seeded on first boot.",
        "AAA Cinemas (Inorbit Mall): 8 Screens (1,552 seats), including Q-Luxon LED & VIP screens.",
        "INOX Varun Beach: 6 Screens, featuring the 37-recliner INSIGNIA auditorium.",
        "Jagadamba 70MM: Iconic 886-seat single-screen auditorium layout.",
        "Mukta A2, Cinepolis, Sri Melody, Sangam Sarat & more."
    ]
    for t in theaters_list:
        p = tf_s5_2.add_paragraph()
        p.text = "✔ " + t
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

    # ==================== SLIDE 6: CONCURRENCY & DOUBLE-BOOKING PREVENTION ====================
    slide6 = add_blank_slide()
    add_header(slide6, "Concurrency Control & Double-Booking Prevention", "Four-layer defensive architecture guaranteeing zero race conditions under high traffic")

    def_layers = [
        ("1. Distributed 8-Minute SeatLock", "Upon seat selection, a SeatLock record with an 8-minute TTL and unique cryptographic token is created. Other users see these seats live as 'Being Booked' via 10-second polling.", CYAN),
        ("2. Database-Level Unique Constraint", "Confirmed seats write to BookedSeat with UNIQUE(showtime_id, seat_id). If two transactions conflict, the database rejects the duplicate at the atomic storage level.", PRIMARY),
        ("3. Atomic Transaction Rollback & Auto-Refund", "If a seat insert fails during gateway callback, the transaction rolls back, marks status as FAILED, and instantly refunds wallet/charges without selling the same seat twice.", ACCENT),
        ("4. Auto-Purge & Lock Expiry Garbage Collection", "Expired pending holds are automatically purged via purge_expired_holds(), ensuring abandoned browser tabs release seats back into the available pool.", SUCCESS)
    ]

    for i, (title, desc, col) in enumerate(def_layers):
        c_top = Inches(1.9 + i * 1.25)
        add_card(slide6, Inches(0.8), c_top, Inches(11.7), Inches(1.1))
        tb = slide6.shapes.add_textbox(Inches(1.1), c_top + Inches(0.12), Inches(11.1), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_LIGHT

    # ==================== SLIDE 7: DYNAMIC PRICING & OFFERS ====================
    slide7 = add_blank_slide()
    add_header(slide7, "Smart Dynamic Pricing & Loyalty Rewards", "Automated revenue management algorithm adapting to real-time seat occupancy and showtime")

    cards_p = [
        ("High Demand Surge", "When screen occupancy reaches >= 60%:\n• Factor: +15% surge price\n• Badge: High Demand Surge\n• Optimizes theater box office revenue during peak blockbusters.", ACCENT),
        ("Early Bird Discount", "When users book >= 2 days in advance:\n• Factor: -10% discount\n• Badge: Early Bird Deal\n• Incentivizes advance bookings and predictable attendance.", CYAN),
        ("Flash Last-Minute Deal", "On show day with < 30% occupancy:\n• Factor: -20% discount\n• Badge: Last-Minute Flash Deal\n• Drives last-minute ticket liquidations for underperforming shows.", SUCCESS),
    ]

    for i, (title, desc, color) in enumerate(cards_p):
        c_left = Inches(0.8 + i * 3.95)
        add_card(slide7, c_left, Inches(1.9), Inches(3.75), Inches(3.2))
        tb = slide7.shapes.add_textbox(c_left + Inches(0.25), Inches(2.1), Inches(3.25), Inches(2.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.space_after = Pt(8)

        for line in desc.split("\n"):
            p_b = tf.add_paragraph()
            p_b.text = line
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = TEXT_LIGHT
            p_b.space_after = Pt(4)

    # Bottom loyalty card
    add_card(slide7, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.4))
    tb_loyalty = slide7.shapes.add_textbox(Inches(1.1), Inches(5.45), Inches(11.1), Inches(1.1))
    tf_l = tb_loyalty.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Loyalty Points & Snacks Add-On Engine"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p2 = tf_l.add_paragraph()
    p2.text = "• Users earn points on registration & bookings; can redeem points directly at checkout (1 pt = ₹1 discount).\n• In-checkout F&B snack bundling (Popcorn, Nachos, Beverages) integrated into the unified checkout invoice."
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_LIGHT

    # ==================== SLIDE 8: PAYMENT GATEWAY & WALLET ====================
    slide8 = add_blank_slide()
    add_header(slide8, "Payment Gateway Integration & Wallet Ledger", "Production-ready Razorpay checkout, sandbox fallback, and in-app credit wallet")

    # Left: Payment Architecture
    add_card(slide8, Inches(0.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb_pay1 = slide8.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf_p1 = tb_pay1.text_frame
    tf_p1.word_wrap = True

    p = tf_p1.paragraphs[0]
    p.text = "Dual Payment Gateway"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(10)

    pay_points = [
        ("Live Razorpay Mode", "Automated switchover when API keys exist. Supports UPI, Cards, Netbanking via Razorpay modal."),
        ("HMAC-SHA256 Verification", "Server validates sha256(order_id + '|' + payment_id, secret). Stops forged client callbacks."),
        ("Sandbox Simulator Mode", "Complete demo mode with simulated success/failure testing, requiring zero third-party accounts."),
        ("Transparent Invoicing", "Base Ticket + ₹20/seat Convenience Fee + 18% GST on fee = Grand Total.")
    ]
    for title, desc in pay_points:
        p = tf_p1.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(6)

    # Right: Wallet System
    add_card(slide8, Inches(6.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb_pay2 = slide8.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf_p2 = tb_pay2.text_frame
    tf_p2.word_wrap = True

    p = tf_p2.paragraphs[0]
    p.text = "CineWave Digital Wallet & Ledger"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.space_after = Pt(10)

    wallet_points = [
        ("₹500 Welcome Credit", "Every registered user automatically receives ₹500 starting wallet balance."),
        ("One-Click Checkout", "Instant debits without third-party redirection when wallet balance covers the total."),
        ("Double-Entry Audit Ledger", "Tracks all transactions (RECHARGE, PAYMENT, REFUND) with timestamps and references."),
        ("Admin Recharge Approval", "Users request wallet recharges; administrators review and approve/reject requests with live credit.")
    ]
    for title, desc in wallet_points:
        p = tf_p2.add_paragraph()
        p.text = f"✔ {title}: {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(6)

    # ==================== SLIDE 9: USER JOURNEY & WORKFLOW ====================
    slide9 = add_blank_slide()
    add_header(slide9, "End-to-End User Experience & Flow", "Seamless 6-step journey from discovery to instant verifiable digital admission")

    steps = [
        ("1. Discovery", "Browse Now Showing & Upcoming films with TMDb posters, ratings & trailers."),
        ("2. Showtime Selection", "Choose preferred Vizag theater, screen format (2D/3D/IMAX), and show timings."),
        ("3. Seat Map Selection", "Select interactive seats across tiers with live price updates and 8-min lock timer."),
        ("4. Checkout & Add-ons", "Add snacks, apply loyalty points, split bills, and review GST invoice."),
        ("5. Payment Execution", "Pay via CineWave Wallet or Razorpay/Sandbox (UPI, Card, Netbanking)."),
        ("6. E-Ticket & QR", "Instant booking confirmation with booking reference and scannable QR pass.")
    ]

    for i, (title, desc) in enumerate(steps):
        col_idx = i % 3
        row_idx = i // 3
        c_left = Inches(0.8 + col_idx * 3.95)
        c_top = Inches(1.9 + row_idx * 2.5)
        
        add_card(slide9, c_left, c_top, Inches(3.75), Inches(2.2))
        tb = slide9.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.2), Inches(3.35), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_LIGHT

    # ==================== SLIDE 10: ADMIN DASHBOARD ====================
    slide10 = add_blank_slide()
    add_header(slide10, "Comprehensive Admin & Cinema Management Hub", "Dedicated portal for theater operations, catalog management, and financial control")

    admin_cols = [
        ("Screen & Layout Builder", "• Add new theaters and auditoriums\n• Interactive seat layout designer\n• Custom row counts, aisles, spans\n• Seed/Reseed all Vizag venues with 1 click", CYAN),
        ("Showtime & Pricing Control", "• Schedule movies to screens & dates\n• Set custom price overrides per class\n• Automated dynamic surge adjustments\n• Import metadata directly via TMDb API", PRIMARY),
        ("Finance & Booking Audits", "• Approve/reject user wallet recharges\n• Real-time booking ledger with status\n• Seat lock monitor & hold management\n• Refund and transaction tracking", ACCENT),
    ]

    for i, (title, desc, color) in enumerate(admin_cols):
        c_left = Inches(0.8 + i * 3.95)
        add_card(slide10, c_left, Inches(1.9), Inches(3.75), Inches(4.8))
        tb = slide10.shapes.add_textbox(c_left + Inches(0.25), Inches(2.2), Inches(3.25), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(17)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.space_after = Pt(12)

        for line in desc.split("\n"):
            p_b = tf.add_paragraph()
            p_b.text = line
            p_b.font.size = Pt(13)
            p_b.font.color.rgb = TEXT_LIGHT
            p_b.space_after = Pt(8)

    # ==================== SLIDE 11: TESTING & SECURITY ====================
    slide11 = add_blank_slide()
    add_header(slide11, "Testing, Verification & Security Architecture", "Robust automated verification suite ensuring bulletproof system reliability")

    # Left: Automated Tests
    add_card(slide11, Inches(0.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb_t1 = slide11.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf_t1 = tb_t1.text_frame
    tf_t1.word_wrap = True

    p = tf_t1.paragraphs[0]
    p.text = "13 Automated End-to-End Tests (test_flow.py)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.space_after = Pt(8)

    test_checks = [
        "1. Seeding of theaters & screens verified",
        "2. Multi-class pricing & layout geometry parsed",
        "3. User authentication & session management",
        "4. 8-minute seat lock creation & token check",
        "5. Concurrency: 2nd user blocked from locked seats",
        "6. Wallet deduction & instant ticket generation",
        "7. Double-booking rejection on confirmed seats",
        "8. Gateway order creation & HMAC verification",
        "9. Failed payment handling & automatic lock purge"
    ]
    for tc in test_checks:
        p = tf_t1.add_paragraph()
        p.text = "✔ " + tc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(4)

    # Right: Security Measures
    add_card(slide11, Inches(6.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb_t2 = slide11.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf_t2 = tb_t2.text_frame
    tf_t2.word_wrap = True

    p = tf_t2.paragraphs[0]
    p.text = "Application Security Mechanisms"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(8)

    sec_points = [
        ("Password Hashing", "Werkzeug salted PBKDF2 hashing; zero plain-text storage."),
        ("Role-Based Access Control", "Decorators protect Admin dashboards and API endpoints."),
        ("Cryptographic Signatures", "HMAC-SHA256 signature validation on payment callbacks."),
        ("ORM SQL Injection Guard", "SQLAlchemy parameterized queries eliminate SQL injection vulnerabilities."),
        ("Tokenized Seat Holds", "Cryptographic session tokens prevent unauthorized hold modifications.")
    ]
    for title, desc in sec_points:
        p = tf_t2.add_paragraph()
        p.text = f"🛡 {title}: {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(6)

    # ==================== SLIDE 12: PROGRESS & MILESTONES ====================
    slide12 = add_blank_slide()
    add_header(slide12, "Milestones Achieved & Roadmap to Final Review", "Review 2 progress assessment and planned final deliverables")

    # Left: Completed in Review 2
    add_card(slide12, Inches(0.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb_m1 = slide12.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf_m1 = tb_m1.text_frame
    tf_m1.word_wrap = True

    p = tf_m1.paragraphs[0]
    p.text = "Completed in Review 2 (Current Status)"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.space_after = Pt(10)

    done_items = [
        "100% Full-Stack MVP Operational",
        "Dynamic JSON Seat Layout Engine for all Vizag halls",
        "Distributed 8-Minute SeatLock concurrency system",
        "Dual Payment Gateway (Razorpay & Sandbox)",
        "In-app Wallet Ledger with Admin Recharge Workflow",
        "Dynamic Surge & Early-Bird Pricing logic",
        "QR Code-enabled digital ticket generation",
        "TMDb Movie Catalog Import integration"
    ]
    for item in done_items:
        p = tf_m1.add_paragraph()
        p.text = "✔ " + item
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(6)

    # Right: Final Review Scope
    add_card(slide12, Inches(6.8), Inches(1.9), Inches(5.7), Inches(4.8))
    tb_m2 = slide12.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf_m2 = tb_m2.text_frame
    tf_m2.word_wrap = True

    p = tf_m2.paragraphs[0]
    p.text = "Planned for Final Review (Final Stage)"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)

    future_items = [
        "Automated Email & SMS Ticket Dispatch (Twilio / SendGrid)",
        "Pre-Showtime Cancellation & Partial Wallet Refund policy",
        "Screen-by-Screen Occupancy Analytics & Revenue Visualizer",
        "Discount Coupon / Voucher Code Engine",
        "PostgreSQL / MySQL Production deployment on Cloud (AWS/Render)",
        "Docker containerization & CI/CD workflow"
    ]
    for item in future_items:
        p = tf_m2.add_paragraph()
        p.text = "➔ " + item
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(6)

    # ==================== SLIDE 13: CONCLUSION & KEY TAKEAWAYS ====================
    slide13 = add_blank_slide()
    add_header(slide13, "Conclusion & Technical Learnings", "Demonstrating full-stack proficiency, distributed thinking, and scalable architecture")

    conclusions = [
        ("Full Stack Mastery", "Seamlessly interconnected Python backend (Flask/SQLAlchemy) with interactive responsive frontend interfaces, RESTful APIs, and external webhooks.", PRIMARY),
        ("Concurrency Engineering", "Solved high-concurrency race conditions using database unique constraints, time-to-live leases, and atomic rollbacks.", CYAN),
        ("Commercial-Grade Features", "Engineered real-world financial checkout workflows, cryptographic payment verification, dynamic surge pricing, and administrative auditing.", SUCCESS),
        ("Production Readiness", "Adhered to secure coding practices (password hashing, signature checking, parameterized queries) backed by end-to-end automated testing.", ACCENT)
    ]

    for i, (title, desc, col) in enumerate(conclusions):
        c_top = Inches(1.9 + i * 1.25)
        add_card(slide13, Inches(0.8), c_top, Inches(11.7), Inches(1.1))
        tb = slide13.shapes.add_textbox(Inches(1.1), c_top + Inches(0.12), Inches(11.1), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_LIGHT

    # ==================== SLIDE 14: Q&A / DEMO ====================
    slide14 = add_blank_slide()
    add_card(slide14, Inches(1.2), Inches(1.2), Inches(10.9), Inches(5.1), CARD_BG, PRIMARY)

    tb_end = slide14.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(9.7), Inches(4.0))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True

    p0 = tf_end.paragraphs[0]
    p0.text = "CINEWAVE CINEMA MANAGEMENT & BOOKING HUB"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = CYAN
    p0.space_after = Pt(16)

    p1 = tf_end.add_paragraph()
    p1.text = "Thank You!"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.space_after = Pt(12)

    p2 = tf_end.add_paragraph()
    p2.text = "Questions & Technical Discussion  |  Ready for Live Demonstration"
    p2.font.size = Pt(16)
    p2.font.color.rgb = PRIMARY
    p2.space_after = Pt(20)

    p3 = tf_end.add_paragraph()
    p3.text = "Demo Highlights: Real Seating Grid • 8-Min Seat Lock • Wallet / Gateway Payment • Instant E-Ticket QR • Admin Layout Editor"
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_MUTED

    output_path = os.path.abspath("CineWave_Project_Review_2_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
